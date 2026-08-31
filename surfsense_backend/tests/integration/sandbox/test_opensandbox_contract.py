"""Live provider contract; opt in with OPENSANDBOX_INTEGRATION=1."""

from __future__ import annotations

import os
import shlex
import subprocess
from mimetypes import guess_type
from pathlib import Path

import pytest

from app.agents.chat.multi_agent_chat.subagents.builtins.deliverables.tools.sandbox import (
    _run_python_script,
)
from app.agents.chat.multi_agent_chat.subagents.builtins.deliverables.tools.save_artifact import (
    _read_artifact_file,
)
from app.artifacts.verification.formats.pdf import check_pdf
from app.artifacts.verification.formats.registry import (
    DOCX_MIME,
    PPTX_MIME,
    get_format_adapter,
)
from app.artifacts.verification.receipt import read_receipt
from app.artifacts.verification.service import verify_artifact
from app.config import config as app_config
from app.sandbox.providers.opensandbox import OpenSandboxProvider

# Matches the tag docker-compose.dev.yml builds for the sandbox-image service.
DEV_SANDBOX_IMAGE = os.getenv("SANDBOX_IMAGE", "surfsense-sandbox:dev")
REPO_ROOT = Path(__file__).resolve().parents[4]
FORMAT_SKILLS = tuple(
    path.parent.name
    for path in sorted((REPO_ROOT / "docker/sandbox/skills").glob("*/SKILL.md"))
)

pytestmark = [
    pytest.mark.integration,
    pytest.mark.skipif(
        os.getenv("OPENSANDBOX_INTEGRATION") != "1",
        reason="requires the local OpenSandbox compose service",
    ),
]


@pytest.mark.parametrize(
    ("skill", "prompt", "expected_mime", "expected_evidence_steps"),
    [
        (
            "pdf",
            "Create a one-page PDF listing three facts about X.",
            "application/pdf",
            ("backend PDF check", "pdftoppm"),
        )
    ],
)
async def test_opensandbox_one_shot_python_binary_io_and_terminate(
    monkeypatch, skill, prompt, expected_mime, expected_evidence_steps
):
    monkeypatch.setattr(app_config, "OPENSANDBOX_DOMAIN", "localhost:8080")
    monkeypatch.setattr(app_config, "OPENSANDBOX_API_KEY", "surfsense-dev-sandbox")
    monkeypatch.setattr(app_config, "SANDBOX_IMAGE", DEV_SANDBOX_IMAGE)
    monkeypatch.setattr(app_config, "SANDBOX_IDLE_TTL_SECONDS", 900)
    provider = OpenSandboxProvider()
    thread_id = "pytest-opensandbox-contract"

    await provider.terminate_session(thread_id)
    session = await provider.get_or_create_session(thread_id)
    try:
        evidence: list[str] = []
        first = await _run_python_script(session, "print(41)")
        second = await _run_python_script(session, "print(42)")
        pdf = await _run_python_script(
            session,
            """
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
pdfmetrics.registerFont(TTFont("DejaVu", "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"))
c = canvas.Canvas("/tmp/three-facts.pdf")
c.setFont("DejaVu", 12)
for y, fact in zip((740, 710, 680), ("Fact one", "Fact two", "Fact three")):
    c.drawString(72, y, fact)
c.save()
""",
        )
        rendered = await session.run_command(
            "mkdir -p /tmp/three-facts-pages && "
            "pdftoppm -jpeg -r 100 /tmp/three-facts.pdf "
            "/tmp/three-facts-pages/page"
        )
        evidence.append("pdftoppm")
        pdf_data = await session.read_file("/tmp/three-facts.pdf")
        checked = check_pdf(pdf_data)
        evidence.insert(0, "backend PDF check")
        jpeg_data = await session.read_file("/tmp/three-facts-pages/page-1.jpg")
        await session.write_file("/tmp/contract.bin", b"\x00SurfSense")
        data = await session.read_file("/tmp/contract.bin")
        stat = await session.run_command("stat -c %s /tmp/contract.bin")

        assert first.ok and "41" in first.output
        assert second.ok and "42" in second.output
        assert pdf.ok and rendered.ok and checked.clean
        assert prompt
        assert guess_type("/tmp/three-facts.pdf")[0] == expected_mime
        assert tuple(evidence) == expected_evidence_steps
        assert pdf_data.startswith(b"%PDF")
        assert jpeg_data.startswith(b"\xff\xd8")
        assert data == b"\x00SurfSense"
        assert int(stat.output.strip()) == len(data)
    finally:
        await provider.terminate_session(thread_id)


@pytest.mark.parametrize(
    ("format_name", "mime_type"),
    [
        ("docx", DOCX_MIME),
        ("pptx", PPTX_MIME),
    ],
)
async def test_live_office_verification_and_mime(monkeypatch, format_name, mime_type):
    monkeypatch.setattr(app_config, "OPENSANDBOX_DOMAIN", "localhost:8080")
    monkeypatch.setattr(app_config, "OPENSANDBOX_API_KEY", "surfsense-dev-sandbox")
    monkeypatch.setattr(app_config, "SANDBOX_IMAGE", DEV_SANDBOX_IMAGE)
    monkeypatch.setattr(app_config, "SANDBOX_IDLE_TTL_SECONDS", 900)
    provider = OpenSandboxProvider()
    thread_id = f"pytest-opensandbox-{format_name}"
    secret = "integration-secret"

    await provider.terminate_session(thread_id)
    session = await provider.get_or_create_session(thread_id)
    try:
        skill_checks = " && ".join(
            f"test -f /opt/skills/{shlex.quote(skill)}/SKILL.md "
            f"&& test ! -d /opt/skills/{shlex.quote(skill)}/scripts"
            for skill in FORMAT_SKILLS
        )
        skills = await session.run_command(skill_checks)
        if format_name == "docx":
            javascript = """
const { Document, Packer, Paragraph, TextRun } = require("docx");
const fs = require("fs");
const doc = new Document({ sections: [{ children: [
  new Paragraph({ children: [new TextRun(
    "A sufficiently long Word document sentence for verification."
  )] }),
  new Paragraph({ children: [new TextRun("email • phone • linkedin")] })
] }] });
Packer.toBuffer(doc).then((buffer) => fs.writeFileSync("/tmp/report.docx", buffer));
"""
            generated = await session.run_command(f"node -e {shlex.quote(javascript)}")
        else:
            generated = await _run_python_script(
                session,
                """
import base64
from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches

presentation = Presentation()
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
slide.shapes.title.text = "Quarterly review"
slide.placeholders[1].text = "A sufficiently descriptive slide for verification."
slide.shapes.add_connector(
    MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(4), Inches(1)
)
slide.shapes.add_connector(
    MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(1), Inches(0), Inches(3)
)
image = BytesIO(base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8"
    "/x8AAusB9Y9Zl1sAAAAASUVORK5CYII="
))
picture = slide.shapes.add_picture(image, Inches(10), Inches(5), width=Inches(1))
picture.crop_left = -0.1
picture.crop_right = 0.2
group = slide.shapes.add_group_shape()
group.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(8), Inches(1), Inches(1), Inches(1)
)
presentation.save("/tmp/report.pptx")
""",
            )
        assert skills.ok and generated.ok

        primary_path = f"/tmp/report.{format_name}"
        result = await verify_artifact(
            session,
            primary_path,
            format=format_name,
            workspace_id=1,
            vision_llm=None,
            secret_key=secret,
        )
        receipt = await read_receipt(
            session,
            secret,
            workspace_id=1,
            primary_path=primary_path,
        )
        stored = await _read_artifact_file(
            session,
            primary_path,
            "primary",
            get_format_adapter(format_name),
        )
        pages = await session.run_command(
            "ls /tmp/surfsense-verify-*/page-*.jpg >/dev/null 2>&1"
        )

        assert result.verified
        assert result.preview_path
        assert receipt.primary_path == primary_path
        assert receipt.preview_path == result.preview_path
        assert not pages.ok
        assert stored.mime_type == mime_type
    finally:
        await provider.terminate_session(thread_id)


def _compose_ip(container: str) -> str:
    probe = subprocess.run(
        [
            "docker",
            "inspect",
            "-f",
            "{{range .NetworkSettings.Networks}}{{.IPAddress}} {{end}}",
            container,
        ],
        capture_output=True,
        text=True,
    )
    addresses = probe.stdout.split()
    if probe.returncode != 0 or not addresses:
        pytest.skip(f"cannot resolve {container} on the compose network")
    return addresses[0]


async def test_sandbox_cannot_reach_the_compose_network(monkeypatch):
    """Egress deny has to hold for raw IPs, not only for hostnames.

    `[egress] mode = "dns"` filters hostnames through a DNS proxy and installs
    no packet filter, so a sandbox could open Postgres at its container IP and
    log in with the default credentials. Only `dns+nft` enforces the policy at
    the network layer, and nothing about the config reads as broken, so the
    connection has to actually be attempted for the regression to be visible.
    """
    monkeypatch.setattr(app_config, "OPENSANDBOX_DOMAIN", "localhost:8080")
    monkeypatch.setattr(app_config, "OPENSANDBOX_API_KEY", "surfsense-dev-sandbox")
    monkeypatch.setattr(app_config, "SANDBOX_IMAGE", DEV_SANDBOX_IMAGE)
    db_ip = _compose_ip(os.getenv("OPENSANDBOX_DB_CONTAINER", "surfsense-dev-db-1"))

    provider = OpenSandboxProvider()
    thread_id = "pytest-opensandbox-egress"
    await provider.terminate_session(thread_id)
    session = await provider.get_or_create_session(thread_id)
    try:
        # Distinguishes "policy blocks it" from "the sandbox runs nothing".
        assert (await session.run_command("echo reachable")).ok

        connect = await session.run_command(
            'python3 -c "import socket,sys;s=socket.socket();s.settimeout(5);'
            f"sys.exit(s.connect_ex(('{db_ip}',5432)))\""
        )
        assert not connect.ok, f"sandbox opened Postgres at {db_ip}:5432"
    finally:
        await provider.terminate_session(thread_id)
